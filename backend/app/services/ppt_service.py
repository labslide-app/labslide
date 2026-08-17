"""PPT 转换服务：提取文本、PPTX/PPT → PDF → PNG。

设计要点：
- 文本提取（.pptx）逐 shape 容错，避免 SmartArt / 特殊形状导致中断。
- LibreOffice 转换使用无空格临时文件名、独立 HOME、超时与进程清理。
- PDF → PNG 使用 pdftoppm，失败时回退 mutool / ImageMagick。
- 输出统一为 1920x1080 PNG（等比缩放 + 白色背景填充）。
"""
import asyncio
import logging
import os
import shutil
import uuid
from pathlib import Path

logger = logging.getLogger(__name__)

# 本地工作目录（转换的中间产物）
UPLOAD_DIR = Path("uploads")
PPTX_DIR = UPLOAD_DIR / "pptx"
PDF_DIR = UPLOAD_DIR / "pdf"
PNG_DIR = UPLOAD_DIR / "png"
WORK_DIR = UPLOAD_DIR / "work"

for d in [PPTX_DIR, PDF_DIR, PNG_DIR, WORK_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# 目标输出尺寸（符合项目硬性约束）
TARGET_WIDTH = 1920
TARGET_HEIGHT = 1080

# 转换超时（秒）
CONVERT_TIMEOUT = 120


# ==================== 工具探测 ====================

def _find_executable(candidates: list[str]) -> str | None:
    """从候选名/路径中查找第一个可用的可执行文件。"""
    for name in candidates:
        # 若是显式路径，直接判断
        if os.path.sep in name or (os.altsep and os.altsep in name):
            if os.path.isfile(name):
                return name
            continue
        # 否则在 PATH 中查找
        found = shutil.which(name)
        if found:
            return found
    return None


def get_libreoffice() -> str | None:
    return _find_executable([
        "soffice",
        "libreoffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
        "/opt/libreoffice/program/soffice",
    ])


def get_pdftoppm() -> str | None:
    return _find_executable([
        "pdftoppm",
        "pdftoppm.exe",
        r"C:\Program Files\Poppler\Library\bin\pdftoppm.exe",
        r"C:\Program Files\poppler\Library\bin\pdftoppm.exe",
        r"C:\Program Files\poppler\bin\pdftoppm.exe",
        "/usr/bin/pdftoppm",
    ])


def get_mutool() -> str | None:
    return _find_executable(["mutool", "mutool.exe"])


def get_imagemagick() -> str | None:
    return _find_executable(["magick", "convert", "convert.exe"])


# ==================== 文本提取 ====================

def _extract_shape_text(shape) -> str:
    """从单个 shape 提取文本，任何异常均返回空字符串。"""
    try:
        if not shape.has_text_frame:
            return ""
        parts = []
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs)
            text = text or paragraph.text
            if text and text.strip():
                parts.append(text.strip())
        return "\n".join(parts)
    except Exception:
        return ""


def extract_text_from_pptx(file_path: str) -> list[dict]:
    """提取每页 PPT 文本与备注（.pptx）。逐 shape / 逐页容错。

    Returns:
        list of dict: [{"page_number": 1, "text": "...", "notes": "..."}, ...]
    """
    from pptx import Presentation as PPTXPresentation

    prs = PPTXPresentation(file_path)
    slides_data = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        try:
            for shape in slide.shapes:
                t = _extract_shape_text(shape)
                if t:
                    texts.append(t)
        except Exception as e:
            logger.warning(f"提取第 {i} 页文本时出错（跳过该页形状）: {e}")

        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                note_parts = []
                for shape in notes_slide.shapes:
                    t = _extract_shape_text(shape)
                    if t:
                        note_parts.append(t)
                notes_text = "\n".join(note_parts)
        except Exception as e:
            logger.warning(f"提取第 {i} 页备注时出错: {e}")

        slides_data.append({
            "page_number": i,
            "text": "\n".join(texts),
            "notes": notes_text,
        })

    return slides_data


# ==================== PPT/PPTX → PDF ====================

async def convert_to_pdf(input_path: str, output_dir: str) -> str:
    """使用 LibreOffice headless 将 PPT/PPTX 转换为 PDF。

    - 复制到工作目录并使用 UUID 文件名，避免空格/特殊字符问题。
    - 设置独立 HOME 到可写目录。
    - 使用 --norestore --nolockcheck 提高稳定性。
    - 超时后终止进程。

    Returns:
        str: 生成的 PDF 文件路径；失败时抛出 RuntimeError。
    """
    soffice = get_libreoffice()
    if not soffice:
        raise RuntimeError("未检测到 LibreOffice，请安装 LibreOffice 以支持 PPT 转换")

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # 复制到工作目录，使用 UUID 作为无特殊字符的文件名
    src = Path(input_path)
    ext = src.suffix or ".pptx"
    safe_name = f"{uuid.uuid4().hex}{ext}"
    work_input = WORK_DIR / safe_name
    shutil.copy2(str(src), str(work_input))

    # 独立的 HOME 目录，避免权限问题
    home_dir = WORK_DIR / "home"
    home_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        soffice,
        "--headless",
        "--norestore",
        "--nolockcheck",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(work_input),
    ]

    env = os.environ.copy()
    env["HOME"] = str(home_dir)
    env["USERPROFILE"] = str(home_dir)

    expected_pdf = output_dir / f"{work_input.stem}.pdf"

    try:
        process = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            env=env,
        )
        try:
            stdout, stderr = await asyncio.wait_for(
                process.communicate(), timeout=CONVERT_TIMEOUT
            )
        except asyncio.TimeoutError:
            process.kill()
            await process.communicate()
            raise RuntimeError(f"LibreOffice 转换超时（{CONVERT_TIMEOUT}s）")

        if process.returncode != 0:
            raise RuntimeError(f"LibreOffice 转换失败: {stderr.decode(errors='ignore')[:500]}")

        if not expected_pdf.exists():
            raise RuntimeError("LibreOffice 未生成 PDF 文件")

        return str(expected_pdf)

    except FileNotFoundError:
        raise RuntimeError(f"LibreOffice 可执行文件无效: {soffice}")
    finally:
        # 清理工作文件
        try:
            if work_input.exists():
                work_input.unlink()
        except Exception:
            pass


# ==================== PDF → PNG ====================

async def _pdftoppm_to_png(pdf_path: str, output_dir: str, prefix: str) -> list[str]:
    """使用 pdftoppm 转换，返回生成的 PNG 路径列表。"""
    pdftoppm = get_pdftoppm()
    if not pdftoppm:
        raise RuntimeError("未检测到 pdftoppm（poppler-utils）")

    cmd = [
        pdftoppm,
        "-png",
        "-r", "150",
        pdf_path,
        str(Path(output_dir) / prefix),
    ]
    process = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    try:
        stdout, stderr = await asyncio.wait_for(process.communicate(), timeout=CONVERT_TIMEOUT)
    except asyncio.TimeoutError:
        process.kill()
        await process.communicate()
        raise RuntimeError("pdftoppm 转换超时")

    if process.returncode != 0:
        raise RuntimeError(f"pdftoppm 转换失败: {stderr.decode(errors='ignore')[:500]}")

    pngs = sorted(
        Path(output_dir).glob(f"{prefix}-*.png"),
        key=lambda p: int(p.stem.rsplit("-", 1)[-1]) if p.stem.rsplit("-", 1)[-1].isdigit() else 0,
    )
    return [str(p) for p in pngs]


async def _mutool_to_png(pdf_path: str, output_dir: str, prefix: str) -> list[str]:
    """使用 mutool 作为 pdftoppm 的后备方案。"""
    mutool = get_mutool()
    if not mutool:
        raise RuntimeError("未检测到 mutool")

    out_pattern = str(Path(output_dir) / f"{prefix}-%d.png")
    cmd = [mutool, "draw", "-o", out_pattern, "-r", "150", pdf_path]
    process = await asyncio.create_subprocess_exec(
        *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    try:
        await asyncio.wait_for(process.communicate(), timeout=CONVERT_TIMEOUT)
    except asyncio.TimeoutError:
        process.kill()
        await process.communicate()
        raise RuntimeError("mutool 转换超时")

    pngs = sorted(
        Path(output_dir).glob(f"{prefix}-*.png"),
        key=lambda p: int(p.stem.rsplit("-", 1)[-1]) if p.stem.rsplit("-", 1)[-1].isdigit() else 0,
    )
    return [str(p) for p in pngs]


async def convert_pdf_to_png(pdf_path: str, output_dir: str) -> list[str]:
    """将 PDF 逐页转换为统一尺寸的 PNG（1920x1080，等比缩放 + 白底填充）。

    Returns:
        list[str]: 按页码排序的 PNG 文件路径列表；全部失败时抛出 RuntimeError。
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    prefix = "slide"

    raw_pngs: list[str] = []
    errors: list[str] = []

    for converter, name in [(_pdftoppm_to_png, "pdftoppm"), (_mutool_to_png, "mutool")]:
        try:
            raw_pngs = await converter(pdf_path, str(output_dir), prefix)
            if raw_pngs:
                break
        except Exception as e:
            errors.append(f"{name}: {e}")
            logger.warning(f"PDF→PNG 转换器 {name} 失败: {e}")

    if not raw_pngs:
        # 尝试 ImageMagick（最后兜底）
        try:
            raw_pngs = await _imagemagick_to_png(pdf_path, str(output_dir), prefix)
        except Exception as e:
            errors.append(f"imagemagick: {e}")

    if not raw_pngs:
        raise RuntimeError("PDF→PNG 转换失败: " + "; ".join(errors))

    # 统一尺寸并填充白底
    normalized = []
    for png in raw_pngs:
        try:
            normalized.append(_normalize_png(png, TARGET_WIDTH, TARGET_HEIGHT))
        except Exception as e:
            logger.warning(f"图片尺寸归一化失败 {png}: {e}")
            normalized.append(png)

    return normalized


async def _imagemagick_to_png(pdf_path: str, output_dir: str, prefix: str) -> list[str]:
    """ImageMagick 兜底方案。"""
    magick = get_imagemagick()
    if not magick:
        raise RuntimeError("未检测到 ImageMagick")

    out_pattern = str(Path(output_dir) / f"{prefix}-%03d.png")
    cmd = [magick, "-density", "150", pdf_path, out_pattern]
    process = await asyncio.create_subprocess_exec(
        *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    try:
        await asyncio.wait_for(process.communicate(), timeout=CONVERT_TIMEOUT)
    except asyncio.TimeoutError:
        process.kill()
        await process.communicate()
        raise RuntimeError("ImageMagick 转换超时")

    pngs = sorted(
        Path(output_dir).glob(f"{prefix}-*.png"),
        key=lambda p: int(p.stem.rsplit("-", 1)[-1]) if p.stem.rsplit("-", 1)[-1].isdigit() else 0,
    )
    return [str(p) for p in pngs]


def _normalize_png(png_path: str, width: int, height: int) -> str:
    """将 PNG 等比缩放到 width x height，并填充白色背景。"""
    from PIL import Image

    img = Image.open(png_path).convert("RGB")
    img.thumbnail((width, height), Image.LANCZOS)

    canvas = Image.new("RGB", (width, height), (255, 255, 255))
    x = (width - img.width) // 2
    y = (height - img.height) // 2
    canvas.paste(img, (x, y))
    canvas.save(png_path, "PNG")
    return png_path


# ==================== 完整处理流程 ====================

def _is_pptx(file_path: str) -> bool:
    return Path(file_path).suffix.lower() == ".pptx"


async def process_presentation(presentation_id: str, file_path: str) -> tuple[list[dict], list[str]]:
    """完整处理 PPT/PPTX：提取文本 → 转 PDF → 转 PNG。

    Returns:
        (slides_data, png_paths): 文本数据与 PNG 路径列表。
    """
    pid = str(presentation_id)

    # 1. 提取文本（.pptx 可提取；.ppt 跳过，后续按图片页数补齐）
    slides_data: list[dict] = []
    if _is_pptx(file_path):
        try:
            slides_data = extract_text_from_pptx(file_path)
        except Exception as e:
            logger.warning(f"PPTX 文本提取失败（继续图片转换）: {e}")
            slides_data = []

    # 2. PPT/PPTX → PDF（convert_to_pdf 返回实际生成的 PDF 路径）
    generated_pdf = await convert_to_pdf(file_path, str(PDF_DIR))
    pdf_path = str(PDF_DIR / f"{pid}.pdf")
    if Path(generated_pdf) != Path(pdf_path):
        shutil.move(generated_pdf, pdf_path)

    if not os.path.exists(pdf_path):
        raise RuntimeError("未找到转换后的 PDF 文件")

    # 3. PDF → PNG
    png_output_dir = str(PNG_DIR / pid)
    png_paths = await convert_pdf_to_png(pdf_path, png_output_dir)
    if not png_paths:
        raise RuntimeError("PDF 转 PNG 无输出")

    # 4. 若 .ppt 未提取文本，按 PNG 数量补齐空文本
    if not slides_data:
        slides_data = [
            {"page_number": i, "text": "", "notes": ""}
            for i in range(1, len(png_paths) + 1)
        ]

    return slides_data, png_paths
